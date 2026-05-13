#!/usr/bin/env python3
"""Build PITCH_DECK.pptx from slide files + brand tokens + PITCH.yaml.

Usage:
    python3 build_pptx.py <slides_dir> <tokens.json> <PITCH.yaml> <output.pptx>

Inputs:
    slides_dir   directory containing slides/NN-*.md with YAML front matter
    tokens.json  brand-assets/tokens.json (colour + type tokens)
    PITCH.yaml   project metadata (company name, stage)
    output.pptx  destination path

Each slide markdown file must begin with YAML front matter declaring at
minimum `headline`, `layout`, `chunk_count`. The body contains the slide
content. A trailing `## Speaker Notes` section becomes embedded notes.

Exit codes:
    0  success
    1  pre-flight failure (missing inputs, malformed front matter)
    2  build failure (python-pptx error during render)
"""

from __future__ import annotations

import json
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path

try:
    import yaml
except ImportError:
    print("ERROR: PyYAML is required. Install with: pip install pyyaml", file=sys.stderr)
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Inches, Pt
except ImportError:
    print(
        "ERROR: python-pptx is required. Install with: pip install python-pptx",
        file=sys.stderr,
    )
    sys.exit(1)


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


@dataclass
class Tokens:
    """Brand tokens consumed by the build."""

    colour_ink: str
    colour_ink_muted: str
    colour_surface: str
    colour_surface_alt: str
    colour_primary: str
    colour_positive: str = "#2e7d32"
    colour_negative: str = "#c62828"
    colour_neutral: str = "#90a4ae"
    font_sans: str = "Calibri"
    font_serif: str = "Georgia"
    text_h1: int = 44
    text_h2: int = 28
    text_lead: int = 22
    text_body: int = 18
    text_caption: int = 13


@dataclass
class Slide:
    """Parsed slide source file."""

    path: Path
    front_matter: dict
    body: str
    speaker_notes: str = ""
    scqa: str = ""
    warnings: list[str] = field(default_factory=list)


def hex_to_rgb(hex_value: str) -> RGBColor:
    """Convert #rrggbb or rrggbb to a python-pptx RGBColor."""
    cleaned = hex_value.lstrip("#")
    if len(cleaned) != 6:
        raise ValueError(f"Invalid hex colour: {hex_value!r}")
    return RGBColor(int(cleaned[0:2], 16), int(cleaned[2:4], 16), int(cleaned[4:6], 16))


def load_tokens(tokens_path: Path) -> Tokens:
    """Read brand-assets/tokens.json into a Tokens dataclass."""
    data = json.loads(tokens_path.read_text())

    # Accept both flat ("colour-ink") and nested ({"colour": {"ink": ...}}) shapes.
    def get(key: str, default: str) -> str:
        if key in data:
            return data[key]
        # Try nested
        parts = key.split("-")
        if len(parts) == 2 and parts[0] in data and isinstance(data[parts[0]], dict):
            return data[parts[0]].get(parts[1], default)
        return default

    return Tokens(
        colour_ink=get("colour-ink", "#1a1a1a"),
        colour_ink_muted=get("colour-ink-muted", "#666666"),
        colour_surface=get("colour-surface", "#ffffff"),
        colour_surface_alt=get("colour-surface-alt", "#f5f5f5"),
        colour_primary=get("colour-primary", "#0066cc"),
        colour_positive=get("colour-positive", "#2e7d32"),
        colour_negative=get("colour-negative", "#c62828"),
        colour_neutral=get("colour-neutral", "#90a4ae"),
        font_sans=get("font-sans", "Calibri"),
        font_serif=get("font-serif", "Georgia"),
        text_h1=int(get("text-h1", "44")),
        text_h2=int(get("text-h2", "28")),
        text_lead=int(get("text-lead", "22")),
        text_body=int(get("text-body", "18")),
        text_caption=int(get("text-caption", "13")),
    )


def parse_slide(path: Path) -> Slide:
    """Parse a slide markdown file with YAML front matter."""
    text = path.read_text()
    match = re.match(r"^---\s*\n(.*?)\n---\s*\n(.*)$", text, re.DOTALL)
    if not match:
        raise ValueError(f"{path.name}: missing YAML front matter")

    front_matter = yaml.safe_load(match.group(1)) or {}
    rest = match.group(2)

    # Split body / speaker notes / SCQA
    parts = re.split(r"^## (Speaker Notes|SCQA.*?)$", rest, flags=re.MULTILINE)
    body = parts[0].strip()
    speaker_notes = ""
    scqa = ""
    for i in range(1, len(parts), 2):
        heading = parts[i]
        content = parts[i + 1].strip() if i + 1 < len(parts) else ""
        if heading.startswith("Speaker"):
            speaker_notes = content
        elif heading.startswith("SCQA"):
            scqa = content

    return Slide(path=path, front_matter=front_matter, body=body, speaker_notes=speaker_notes, scqa=scqa)


def contrast_ratio(rgb1: tuple[int, int, int], rgb2: tuple[int, int, int]) -> float:
    """WCAG relative-luminance contrast ratio between two RGB triples."""

    def relative_luminance(rgb: tuple[int, int, int]) -> float:
        def channel(c: int) -> float:
            v = c / 255.0
            return v / 12.92 if v <= 0.03928 else ((v + 0.055) / 1.055) ** 2.4

        r, g, b = (channel(c) for c in rgb)
        return 0.2126 * r + 0.7152 * g + 0.0722 * b

    l1, l2 = relative_luminance(rgb1), relative_luminance(rgb2)
    if l1 < l2:
        l1, l2 = l2, l1
    return (l1 + 0.05) / (l2 + 0.05)


def hex_to_tuple(hex_value: str) -> tuple[int, int, int]:
    cleaned = hex_value.lstrip("#")
    return int(cleaned[0:2], 16), int(cleaned[2:4], 16), int(cleaned[4:6], 16)


def add_text_box(slide_obj, text: str, left, top, width, height, *, font, size, bold=False, colour=None, align=PP_ALIGN.LEFT):
    """Add a text box with the given text and style."""
    box = slide_obj.shapes.add_textbox(left, top, width, height)
    text_frame = box.text_frame
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    if colour is not None:
        run.font.color.rgb = colour
    return box


def add_background(slide_obj, colour: RGBColor) -> None:
    """Add a full-bleed background rectangle in the given colour."""
    rect = slide_obj.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    rect.line.fill.background()
    rect.fill.solid()
    rect.fill.fore_color.rgb = colour
    # Send to back
    spTree = rect._element.getparent()
    spTree.remove(rect._element)
    spTree.insert(2, rect._element)


def render_slide(prs, slide_obj_data: Slide, tokens: Tokens, pitch_meta: dict) -> list[str]:
    """Render one slide. Returns a list of build warnings."""
    warnings: list[str] = []

    layout = prs.slide_layouts[6]  # blank layout
    slide_obj = prs.slides.add_slide(layout)

    add_background(slide_obj, hex_to_rgb(tokens.colour_surface))

    fm = slide_obj_data.front_matter
    headline = fm.get("headline", "[MISSING HEADLINE]")
    layout_kind = fm.get("layout", "title-content")
    chunk_count = fm.get("chunk_count", 0)
    proof_points = fm.get("proof_points", []) or []

    # Pre-flight warnings
    if chunk_count and chunk_count > 5:
        warnings.append(f"{slide_obj_data.path.name}: chunk_count={chunk_count} exceeds VD-02 limit of 5")
    if not headline or headline.endswith(":") or len(headline.split()) <= 2:
        warnings.append(f"{slide_obj_data.path.name}: headline may be label, not conclusion (ND-04)")

    # Headline
    add_text_box(
        slide_obj,
        headline,
        Inches(0.6),
        Inches(0.5),
        Inches(12.1),
        Inches(1.3),
        font=tokens.font_sans,
        size=tokens.text_h1,
        bold=True,
        colour=hex_to_rgb(tokens.colour_ink),
    )

    # Body (lines from markdown bullets, stripped)
    body_lines = []
    for line in slide_obj_data.body.split("\n"):
        line = line.strip()
        if line.startswith("- ") or line.startswith("* "):
            body_lines.append(line[2:].strip())
        elif line and not line.startswith("#") and not line.startswith("<!--"):
            body_lines.append(line)

    if layout_kind == "title":
        # Centred big headline only; body is optional subtitle
        if body_lines:
            add_text_box(
                slide_obj,
                " ".join(body_lines),
                Inches(0.6),
                Inches(3.5),
                Inches(12.1),
                Inches(2.0),
                font=tokens.font_sans,
                size=tokens.text_lead,
                colour=hex_to_rgb(tokens.colour_ink_muted),
                align=PP_ALIGN.CENTER,
            )
    elif layout_kind in ("title-content", "quote", "team"):
        # Body as a stacked list
        top = Inches(2.2)
        body_box = slide_obj.shapes.add_textbox(Inches(0.6), top, Inches(12.1), Inches(4.5))
        text_frame = body_box.text_frame
        text_frame.word_wrap = True
        for idx, line in enumerate(body_lines):
            paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(8)
            run = paragraph.add_run()
            run.text = f"•  {line}" if layout_kind == "title-content" else line
            run.font.name = tokens.font_sans
            run.font.size = Pt(tokens.text_body)
            run.font.color.rgb = hex_to_rgb(tokens.colour_ink)
    elif layout_kind == "two-column":
        # Body lines split between two columns
        midpoint = (len(body_lines) + 1) // 2
        left_lines = body_lines[:midpoint]
        right_lines = body_lines[midpoint:]
        for col_idx, lines in enumerate((left_lines, right_lines)):
            left = Inches(0.6 if col_idx == 0 else 6.95)
            box = slide_obj.shapes.add_textbox(left, Inches(2.2), Inches(5.8), Inches(4.5))
            tf = box.text_frame
            tf.word_wrap = True
            for idx, line in enumerate(lines):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(8)
                r = p.add_run()
                r.text = f"•  {line}"
                r.font.name = tokens.font_sans
                r.font.size = Pt(tokens.text_body)
                r.font.color.rgb = hex_to_rgb(tokens.colour_ink)
    else:
        # chart, full-image, comparison-matrix — placeholder rendering with a note
        add_text_box(
            slide_obj,
            f"[Layout '{layout_kind}' renders as content area in v0.1 — refine in PPTX manually]",
            Inches(0.6),
            Inches(3.5),
            Inches(12.1),
            Inches(2.0),
            font=tokens.font_sans,
            size=tokens.text_caption,
            colour=hex_to_rgb(tokens.colour_ink_muted),
            align=PP_ALIGN.CENTER,
        )
        if body_lines:
            add_text_box(
                slide_obj,
                " · ".join(body_lines),
                Inches(0.6),
                Inches(4.5),
                Inches(12.1),
                Inches(2.0),
                font=tokens.font_sans,
                size=tokens.text_body,
                colour=hex_to_rgb(tokens.colour_ink),
                align=PP_ALIGN.CENTER,
            )

    # Footer with proof-points
    if proof_points:
        footer = "Sources: " + ", ".join(proof_points)
        add_text_box(
            slide_obj,
            footer,
            Inches(0.6),
            Inches(7.0),
            Inches(12.1),
            Inches(0.3),
            font=tokens.font_sans,
            size=tokens.text_caption,
            colour=hex_to_rgb(tokens.colour_ink_muted),
        )

    # Slide number + company name
    company_name = pitch_meta.get("name", "")
    add_text_box(
        slide_obj,
        company_name,
        Inches(0.6),
        Inches(7.0),
        Inches(6.0),
        Inches(0.3),
        font=tokens.font_sans,
        size=tokens.text_caption,
        colour=hex_to_rgb(tokens.colour_ink_muted),
    )

    # Contrast check
    ratio = contrast_ratio(hex_to_tuple(tokens.colour_ink), hex_to_tuple(tokens.colour_surface))
    if ratio < 4.5:
        warnings.append(
            f"{slide_obj_data.path.name}: ink-on-surface contrast {ratio:.2f}:1 below WCAG AA 4.5:1 (VD-09)"
        )

    # Speaker notes
    if slide_obj_data.speaker_notes:
        notes_frame = slide_obj.notes_slide.notes_text_frame
        notes_frame.text = slide_obj_data.speaker_notes

    return warnings


def build(slides_dir: Path, tokens_path: Path, pitch_path: Path, output_path: Path) -> int:
    if not slides_dir.is_dir():
        print(f"ERROR: slides_dir not found: {slides_dir}", file=sys.stderr)
        return 1
    if not tokens_path.is_file():
        print(f"ERROR: tokens.json not found: {tokens_path}", file=sys.stderr)
        return 1
    if not pitch_path.is_file():
        print(f"ERROR: PITCH.yaml not found: {pitch_path}", file=sys.stderr)
        return 1

    tokens = load_tokens(tokens_path)
    pitch_meta = yaml.safe_load(pitch_path.read_text()) or {}

    slide_files = sorted(slides_dir.glob("*.md"))
    if not slide_files:
        print(f"ERROR: no slide files found in {slides_dir}", file=sys.stderr)
        return 1

    slides = []
    for path in slide_files:
        try:
            slides.append(parse_slide(path))
        except ValueError as exc:
            print(f"ERROR: {exc}", file=sys.stderr)
            return 1

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    all_warnings: list[str] = []
    for slide in slides:
        try:
            warnings = render_slide(prs, slide, tokens, pitch_meta)
            all_warnings.extend(warnings)
        except Exception as exc:  # noqa: BLE001
            print(f"ERROR rendering {slide.path.name}: {exc}", file=sys.stderr)
            return 2

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    if all_warnings:
        print("BUILD WARNINGS:", file=sys.stderr)
        for w in all_warnings:
            print(f"  - {w}", file=sys.stderr)

    print(f"Wrote {output_path} ({len(slides)} slides)")
    return 0


def main(argv: list[str]) -> int:
    if len(argv) != 5:
        print(__doc__, file=sys.stderr)
        return 1
    return build(Path(argv[1]), Path(argv[2]), Path(argv[3]), Path(argv[4]))


if __name__ == "__main__":
    sys.exit(main(sys.argv))
